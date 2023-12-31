import requests
import random
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import math
from pptx.dml.color import RGBColor
import io
import urllib.request
import sys
from PyQt5.QtWidgets import (QApplication, QWidget, QVBoxLayout, QLabel, QLineEdit, QPushButton, QComboBox, QMessageBox)
import uuid
import re

# Define your API keys (replace with your actual keys)
UNSPLASH_API_KEY = 'unsplash_key_goes_here'
OPENAI_API_KEY = 'openai_key_goes_here'


# Define these at the top of your script, after your imports
MAX_BULLETS = 3  # Limit the number of bullet points
MAX_WORDS_PER_BULLET = 10  # Limit the number of words per bullet point
MAX_IMAGES_PER_SLIDE = 1


def truncate_sentences(text, max_sentences=2):
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
    return ' '.join(sentences[:max_sentences])


# Function to fetch a random image URL from Unsplash
def fetch_random_image_url(api_key, search_term='presentation'):
    headers = {
        "Accept-Version": "v1",
        "Authorization": f"Client-ID {api_key}"
    }
    params = {
        'query': search_term,
        'page': 1,
        'per_page': 30
    }
    response = requests.get('https://api.unsplash.com/search/photos', headers=headers, params=params)
    if response.status_code == 200:
        json_response = response.json()
        results = json_response.get('results', [])
        if results:
            random_image = random.choice(results)
            return random_image['urls']['regular']
        else:
            print("No images found for the given search term.")
            return None
    else:
        print(f"Failed to fetch images: Status code {response.status_code}")
        return None

def generate_slide_title(api_key, prompt, slide_number, previous_titles):
    openai.api_key = api_key

    try:
        # Contextual prompt for logical progression without slide numbers
        context = " ".join(previous_titles[-3:])  # Use only the last 3 titles for context to avoid prompt getting too long
        adjusted_prompt = f"Create a concise title that introduces a new aspect of '{prompt}', without numbering, considering previous topics: {context}."

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "system", "content": "You are a helpful assistant."}, 
                      {"role": "user", "content": adjusted_prompt}]
        )

        title = response['choices'][0]['message']['content'].strip().strip('"')

        # Remove any slide numbering added by the model
        title = re.sub(r'^Slide \d+:?\s*', '', title, flags=re.I)

        # Limit title length
        max_title_length = 70  # or any number you see fit
        if len(title) > max_title_length:
            title = title[:max_title_length].rsplit(' ', 1)[0] + '...'  # Avoid cutting words in half

        if title in previous_titles:
            print(f"Warning: Duplicate or empty title for slide {slide_number}. Generated title: '{title}'")
            title = f"{prompt} Aspect {slide_number}"  # Fallback title

        return title

    except Exception as e:
        print(f"An error occurred while generating slide title: {e}")
        return f"{prompt} Aspect {slide_number}"  # Fallback title



def generate_unique_slide_content(api_key, slide_title, prompt, max_bullets=4, max_words_per_bullet=20):
    openai.api_key = api_key

    try:
        adjusted_prompt = f"Provide a concise summary in {max_bullets} bullet points, each with no more than {max_words_per_bullet} words, explaining the key aspects of '{slide_title}' relevant to the theme '{prompt}'."

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "system", "content": "You are a helpful assistant."}, 
                      {"role": "user", "content": adjusted_prompt}]
        )

        content = response['choices'][0]['message']['content']
        bullets = [line.strip() for line in content.split('\n') if line.strip().startswith('-')]
        # Trim each bullet to the word limit
        trimmed_bullets = []
        for bullet in bullets:
            words = bullet.strip('- ').split()
            trimmed_bullet = ' '.join(words[:max_words_per_bullet])
            trimmed_bullets.append(trimmed_bullet)
        
        return trimmed_bullets[:max_bullets]

    except Exception as e:
        print(f"An error occurred: {e}")
        return [f"Details on '{slide_title}' will be discussed."]


    
def estimate_text_height(text, font_size_pt, slide_width_px, slide_height_px, text_box_margin_px=50):
    # Assume an average character width at 12pt font size and scale according to actual font size
    avg_char_width_pt = (Pt(12) / 2.0) * (font_size_pt / Pt(12))
    slide_width_pt = Pt(slide_width_px)
    slide_height_pt = Pt(slide_height_px)
    max_chars_per_line = int((slide_width_pt - Pt(text_box_margin_px*2)) / avg_char_width_pt)
    
    line_count = 0
    for paragraph in text.split('\n'):
        line_count += math.ceil(len(paragraph) / max_chars_per_line)
    
    # Approximate line height is 1.2 times the font size
    total_text_height_pt = line_count * (font_size_pt * 1.2)
    return total_text_height_pt

def create_presentation(prompt, num_slides, api_key_unsplash, api_key_openai):
    prs = Presentation()
    previous_titles = []
    previous_contents = []  # Initialize the variable to store all generated content

    for slide_number in range(1, num_slides + 1):
        slide_layout = prs.slide_layouts[5]  # Use a blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Generate the slide title without slide numbering
        slide_title = generate_slide_title(api_key_openai, prompt, slide_number, previous_titles).replace(f"Slide {slide_number}:", "").strip()
        previous_titles.append(slide_title)

        # Fetch an image relevant to the slide title
        search_term = f"{slide_title} Doctor Who"  # Customize your search term based on slide content
        image_url = fetch_random_image_url(api_key_unsplash, search_term)

        # Generate bullet points with a controlled length
        bullets = generate_unique_slide_content(api_key_openai, slide_title, prompt, max_bullets=4, max_words_per_bullet=20)


        previous_contents.append(' '.join(bullets))  # Track all generated bullet points


        if bullets is None:
            bullets = []  # Ensure bullets is a list to prevent iteration errors

        # Add title to the slide
        title_shape = slide.shapes.title
        title_shape.text = slide_title

        # Define text box position and size
        left = Inches(0.5)
        top = Inches(1.5)
        text_width = Inches(5)  # Half the slide width for text
        text_height = Inches(4.5)  # Adjust as needed

        # Add text box
        textbox = slide.shapes.add_textbox(left, top, text_width, text_height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Add bullet points to text box
        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

        # Check if text overflows and adjust font size if necessary
        if len(text_frame.paragraphs) > 6:  # Threshold for text overflow
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(16)  # Reduce font size

        # Update previous_contents to include the generated content
        slide_content = f"{slide_title}\n{' '.join(bullets)}"
        previous_contents.append(slide_content)

        # Fetch and add an image beside the text box
        for img_num in range(MAX_IMAGES_PER_SLIDE):
            random_image_url = fetch_random_image_url(api_key_unsplash, prompt)
            if random_image_url:
                response = urllib.request.urlopen(random_image_url)
                image_stream = io.BytesIO(response.read())
                image_left = left + text_width + Inches(0.5)  # Space between text and image
                image_top = top
                image_width = prs.slide_width - text_width - Inches(1.0)  # Remaining width
                image_height = prs.slide_height - top - Inches(1.0)  # Remaining height
                img_shape = slide.shapes.add_picture(image_stream, image_left, image_top, image_width, image_height)

    # Generate a file name based on the prompt or a random UUID
    if prompt and all(c.isalnum() or c.isspace() for c in prompt):
        # Use a sanitized version of the prompt as the file name
        file_name = re.sub(r'\W+', '', prompt) + '.pptx'
    else:
        # Generate a random file name
        file_name = str(uuid.uuid4()) + '.pptx'

    # Save the presentation
    prs.save(file_name)
    print(f"Presentation saved as: {file_name}")



class SlideGeneratorApp(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        self.setWindowTitle('PowerPoint Slide Generator')
        self.setGeometry(300, 300, 400, 200)  # x, y, width, height

        layout = QVBoxLayout(self)

        # Topic label and entry
        label_topic = QLabel('Enter Topic:', self)
        self.entry_topic = QLineEdit(self)

        # Number of slides label and combo box
        label_slides = QLabel('Select Number of Slides:', self)
        self.combo_slides = QComboBox(self)
        self.combo_slides.addItems(['1', '2', '3', '4', '5'])

        # Generate button
        self.generate_button = QPushButton('Generate Presentation', self)
        self.generate_button.clicked.connect(self.on_generate_button_click)

        # Add widgets to layout
        layout.addWidget(label_topic)
        layout.addWidget(self.entry_topic)
        layout.addWidget(label_slides)
        layout.addWidget(self.combo_slides)
        layout.addWidget(self.generate_button)

        # Set the layout on the application's window
        self.setLayout(layout)

    def on_generate_button_click(self):
        topic = self.entry_topic.text()
        num_slides = int(self.combo_slides.currentText())
        try:
            # Now passing the required API keys to the function
            create_presentation(topic, num_slides, UNSPLASH_API_KEY, OPENAI_API_KEY)
            QMessageBox.information(self, 'Success', 'Presentation generated successfully.')
        except ValueError as e:
            QMessageBox.critical(self, 'Error', 'Please enter a valid number for the number of slides.')
        except Exception as e:
            QMessageBox.critical(self, 'Error', f'An error occurred: {e}')


# Assume the rest of your functions (fetch_random_image_url, generate_slide_content, create_presentation) are defined here

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = SlideGeneratorApp()
    ex.show()
    sys.exit(app.exec_())
